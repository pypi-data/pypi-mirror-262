import docx
import pandas as pd
import json
from pypdf import PdfReader 
from bs4 import BeautifulSoup
from pptx import Presentation
import codecs
import os

class FileFuse:
    def read(file_name):
        if not Helper.file_exists(file_name):
           raise Exception("File Not Found")
        if len(file_name.split('.')) <=1 :
            raise Exception("File Extension not Found")
        
        extension = file_name.split('.')[-1]
        if extension == 'txt':
            return Helper.read_txt(file_name)
        elif extension == 'docx' or extension == 'doc':
            return Helper.read_docx(file_name)
        elif extension == 'csv':
            return Helper.read_csv(file_name)
        elif extension == 'json':
            return Helper.read_json(file_name)
        elif extension == 'pdf':
            return Helper.read_pdf(file_name)
        elif extension == 'html':
            return Helper.read_html(file_name)
        elif extension == 'xlsx':
            return Helper.read_excel(file_name)
        elif extension == 'pptx':
            return Helper.read_pptx(file_name)
        elif extension == "css" or extension == "js" or extension == "py" or extension == "ts" or extension == "java" or extension == "cpp" or exteions == "c" or extension == "php" or extension == 'rb' or extension == "swift" or extension == "go" or extension == "pl" or extension == "cs" or extension == "jsx" or extension == "scss" or extension == "less" or extension == "vue" or extension == "xml" or extension == "sql" or extension == "xml" or extension == "md" or extension == "yml" or extension == "yaml" or extension == "ini" or extension == "sh" or extension == "env" :
            return Helper.read_code(file_name)
        else:
            try:
                return Helper.read_code(file_name)
            except:
                return "Not Supported + " + str(file_name)


class Helper: 
    def read_html(file_path):
        data = ""
        with open(file_path, 'r') as file:
            soup = BeautifulSoup(file, 'html.parser')
            data = soup.get_text()
        return data
                
    def read_txt(file_path):
        data = ''
        with open(file_path, 'r') as file:
                data = file.read()
        return data

    def read_docx(file_path):
        data = ''
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        data = '\n'.join(full_text)
        return data

    def read_csv(file_path):
        df = pd.read_csv(file_path)
        data = df.to_string()
        return data

    def read_pdf(file_path):
        data = ''
        reader = PdfReader(file_path)
        for x in range(0,len(reader.pages)):
            data = data +"\n\n" + reader.pages[x].extract_text() 
        return data

    def read_json(file_path):
        data = ''
        with open(file_path, 'r') as file:
            data = json.load(file)
        data = json.dumps(data)
        return data

    def read_excel(file_path):
        df = pd.read_excel(file_path)
        data = df.to_string()
        return data

    def read_pptx(file_path):
        pres = Presentation(file_path)
        full_text = []
        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            full_text.append(run.text)
        return ' '.join(full_text)

    def read_code(file_path):
        with open(file_path, 'r') as file:
            data = file.read()
        return data

    def file_exists(file_path):
        if os.path.exists(file_path):
            return True
        return False
